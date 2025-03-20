from pptx import Presentation

# プレゼンテーションを読み込む
prs = Presentation('research_presentation.pptx')

# 各スライドの内容を表示
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            print(shape.text)
