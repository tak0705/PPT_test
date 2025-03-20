from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# フォント変更のモジュール
def set_title_font(title_shape, size, bold, color, font_name):
    """
    タイトルのフォントを設定する関数
    title_shape: タイトルのshapeオブジェクト
    size: フォントサイズ
    bold: 太字設定
    color: フォントカラー
    font_name: フォント名
    """
    title_format = title_shape.text_frame.paragraphs[0].font
    title_format.size = Pt(size)
    title_format.bold = bold
    title_format.color.rgb = color
    title_format.name = font_name  # フォント名をメイリオに設定

# コンテンツフォント変更のモジュール
def set_content_font(content_shape, size, bold, color=RGBColor(0, 0, 0), font_name='メイリオ'):
    """
    コンテンツのフォントを設定する関数
    content_shape: コンテンツのshapeオブジェクト
    size: フォントサイズ
    bold: 太字設定
    color: フォントカラー
    font_name: フォント名
    """
    # 各段落に対してフォント設定
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = font_name  # フォント名を設定

# 背景色の変更のモジュール
def set_slide_background(slide, color=RGBColor(255, 255, 255)):
    """
    スライドの背景色を変更する関数
    slide: スライドオブジェクト
    color: 背景色
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# タイトルの下に青い線を引くモジュール
def add_blue_underline(slide, title_shape):
    """
    タイトルの下に青い線を引く関数
    slide: スライドオブジェクト
    title_shape: タイトルのshapeオブジェクト
    """
    # タイトルの位置を取得
    left = title_shape.left
    top = title_shape.top + title_shape.height
    width = title_shape.width
    height = Inches(0.1)  # 線の太さを0.1インチに設定
    
    # 青い線を描画
    line = slide.shapes.add_shape(
        1,  # msoShapeRectangle (長方形)
        left, top, width, height
    )
    fill = line.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # 青色に設定

    # 線の枠線を消す
    line.line.fill.background()

# プレゼンテーション作成
prs = Presentation()

# --- スライド1: タイトル ---
slide_layout = prs.slide_layouts[0]  # タイトルスライド
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "XX技術の進化とその応用"
subtitle.text = "山田 太郎\n2025年3月20日"

# タイトルフォント設定
set_title_font(title, size=40, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(subtitle, size=18, bold=False, color=RGBColor(102, 102, 102), font_name='メイリオ')

# 背景色設定
set_slide_background(slide, color=RGBColor(240, 240, 240))  # 背景を淡いグレーに

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド2: 研究の目的 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "研究の目的"
content.text = ("XX技術の現状と課題\n"
                "技術の進化に向けたニーズ\n"
                "本研究の目的は、XX技術の応用を拡大すること")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド3: 研究の動機・課題設定 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "研究の動機・課題設定"

# 画像をスライドに追加
img_path = "path_to_image.png"  # 画像ファイルのパスを指定
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(6), height=Inches(3.5))

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド4: 研究方法 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "研究方法"
content.text = ("実験装置の構成\n"
                "実施した実験の手順\n"
                "データ収集方法")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド5: 研究結果 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "研究結果"
content.text = ("実験の結果、XX技術の性能向上が確認されました。\n"
                "特に、YYの面で顕著な改善が見られました。")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド6: グラフや図の挿入 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "研究結果 - グラフ"

# グラフをスライドに追加（画像として）
img_path = "path_to_image.png"  # グラフ画像のパスを指定
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(6), height=Inches(3.5))

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド7: 結論 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "結論"
content.text = ("本研究の要点\n"
                "XX技術の重要性\n"
                "研究の意義と社会への貢献")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド8: 今後の課題 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "今後の課題"
content.text = ("今後、XX技術の実用化に向けた課題として\n"
                "実験結果の精度向上\n"
                "実験条件の最適化\n"
                "技術の拡張が考えられます。")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド9: まとめ ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "まとめ"
content.text = ("XX技術は今後の技術革新に大きく貢献する可能性を秘めており、\n"
                "本研究の成果がその発展に寄与することを期待している。")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# --- スライド10: 参考文献 ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "参考文献"
content.text = ("研究論文A, 2023年\n"
                "研究論文B, 2024年\n"
                "研究論文C, 2025年")

set_title_font(title, size=35, bold=True, color=RGBColor(0, 51, 102), font_name='メイリオ')
set_content_font(content, size=30, bold=False, color=RGBColor(0, 0, 0), font_name='メイリオ')
set_slide_background(slide, color=RGBColor(255, 255, 255))  # 白色背景

# タイトル下に青い傍線を追加
add_blue_underline(slide, title)

# 保存
prs.save('presentation.pptx')
print("PowerPointファイル 'presentation.pptx' が作成されました。")

